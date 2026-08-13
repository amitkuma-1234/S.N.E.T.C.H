# ============================================================
#  document_chatbot.py — Document AI Chatbot (S.N.E.T.C.H)
#  Groq API (qwen/qwen3.6-27b) + hybrid RAG (semantic + keyword) over documents
# ============================================================

import io
import os
import re
import json
import time
import uuid
import queue
import sqlite3
import logging
import threading
import urllib.request
import urllib.error
from datetime import datetime
from typing import Callable

from flask import (
    Blueprint, request, jsonify, Response, stream_with_context, send_file,
)
from werkzeug.utils import secure_filename

from dotenv import load_dotenv
load_dotenv()

def speak(text: str) -> None:
    """Text-only output (was TTS)."""
    print(f"[SNETCH] {text}")

try:
    from utils.logger import logger
except ImportError:
    logger = logging.getLogger("document_chatbot")
    logging.basicConfig(level=logging.INFO)

logging.getLogger("httpx").setLevel(logging.WARNING)
logging.getLogger("httpcore").setLevel(logging.WARNING)

try:
    import numpy as np
    from sentence_transformers import SentenceTransformer, CrossEncoder
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
    DOC_ML_OK = True
except ImportError:
    DOC_ML_OK = False
    logger.warning(
        "Document chatbot ML deps missing. Run: "
        "pip install sentence-transformers scikit-learn numpy"
    )

try:
    import fitz  # PyMuPDF
    PDF_OK = True
except ImportError:
    PDF_OK = False

try:
    from docx import Document as DocxDocument
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

try:
    from pptx import Presentation as PptxPresentation
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

DOC_DEPS_OK = DOC_ML_OK and (PDF_OK or DOCX_OK)

try:
    from langchain_groq import ChatGroq
    from langchain_core.messages import HumanMessage
    from langgraph.graph import StateGraph, START, END
    from langgraph.graph.message import add_messages
    from langgraph.checkpoint.memory import MemorySaver
    from typing import TypedDict, Annotated
    LANGGRAPH_OK = True
except ImportError:
    LANGGRAPH_OK = False

try:
    from googlesearch import search as google_search
    WEB_SEARCH_OK = True
except ImportError:
    WEB_SEARCH_OK = False

# ── Config ──────────────────────────────────────────────────
GROQ_API_KEY       = os.getenv("GROQ_API_KEY", "")
GROQ_CHAT_URL      = "https://api.groq.com/openai/v1/chat/completions"
GROQ_MODEL         = os.getenv("GROQ_MODEL", "qwen/qwen3.6-27b")
OLLAMA_TIMEOUT_SEC = 120
NOT_IN_DOCUMENT    = "This information is not available in the document."

CHUNK_MAX_WORDS    = 120
CHUNK_OVERLAP_SENTS = 2
FULL_DOCUMENT_WORD_LIMIT = 3500

_embed_model = None
_reranker = None
_chat_model = None
_graph_store = {}
_sessions = {}
_active_doc_id = None

# ── Web feature config (Flask blueprint, SQLite, per-thread sessions) ──
BASE_DIR      = os.path.dirname(os.path.abspath(__file__))
DB_STORAGE_DIR = os.path.join(BASE_DIR, "db_storage")
DOC_DB_PATH    = os.path.join(DB_STORAGE_DIR, "document_chatbot.db")
CHECKPOINT_DB_PATH = os.path.join(DB_STORAGE_DIR, "document_chatbot_checkpoints.db")
UPLOAD_DIR     = os.path.join(DB_STORAGE_DIR, "doc_chat_uploads")

ALLOWED_EXTENSIONS = {"pdf", "docx", "doc", "txt", "text", "md", "markdown", "pptx", "ppt"}
MAX_UPLOAD_BYTES   = 50 * 1024 * 1024  # 50 MB

_web_sessions = {}      # thread_id -> session dict (same shape as _new_session)
_web_graph_store = {}   # thread_id -> {"graph":..., "config":...}
_checkpointer_singleton = None


# ════════════════════════════════════════════════════════════
#  OLLAMA
# ════════════════════════════════════════════════════════════

def is_ollama_ready() -> bool:
    return bool(GROQ_API_KEY)


def _ollama_error() -> str:
    return "Groq API key missing or invalid. Add GROQ_API_KEY to your .env file."


def ask_ollama(prompt: str, temperature: float = 0.2, system: str = "") -> str:
    if not GROQ_API_KEY:
        raise RuntimeError(_ollama_error())
    messages = []
    if system:
        messages.append({"role": "system", "content": system})
    messages.append({"role": "user", "content": prompt})

    payload = json.dumps({
        "model": GROQ_MODEL,
        "messages": messages,
        "stream": False,
        "temperature": temperature,
    }).encode("utf-8")

    req = urllib.request.Request(
        GROQ_CHAT_URL, data=payload,
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Bearer {GROQ_API_KEY}",
        }, method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=OLLAMA_TIMEOUT_SEC) as resp:
            body = json.loads(resp.read().decode("utf-8"))
            return body["choices"][0]["message"]["content"].strip()
    except urllib.error.HTTPError as e:
        raise RuntimeError(f"Groq API error {e.code}: {e.read().decode('utf-8', 'ignore')}")
    except urllib.error.URLError as e:
        raise RuntimeError(f"Groq API not reachable — {e}")
    except Exception as e:
        raise RuntimeError(f"Groq call failed: {e}")


def ask_ollama_stream(prompt: str, temperature: float = 0.2, system: str = "",
                      on_token: Callable[[str], None] | None = None) -> str:
    if not GROQ_API_KEY:
        raise RuntimeError(_ollama_error())
    messages = []
    if system:
        messages.append({"role": "system", "content": system})
    messages.append({"role": "user", "content": prompt})

    payload = json.dumps({
        "model": GROQ_MODEL,
        "messages": messages,
        "stream": True,
        "temperature": temperature,
    }).encode("utf-8")

    req = urllib.request.Request(
        GROQ_CHAT_URL, data=payload,
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Bearer {GROQ_API_KEY}",
        }, method="POST",
    )
    parts = []
    try:
        with urllib.request.urlopen(req, timeout=OLLAMA_TIMEOUT_SEC) as resp:
            for raw_line in resp:
                line = raw_line.decode("utf-8").strip()
                if not line or not line.startswith("data:"):
                    continue
                data_str = line[len("data:"):].strip()
                if data_str == "[DONE]":
                    break
                chunk = json.loads(data_str)
                delta = chunk.get("choices", [{}])[0].get("delta", {})
                token = delta.get("content", "")
                if token:
                    parts.append(token)
                    if on_token:
                        on_token(token)
    except urllib.error.HTTPError as e:
        raise RuntimeError(f"Groq API error {e.code}: {e.read().decode('utf-8', 'ignore')}")
    except urllib.error.URLError as e:
        raise RuntimeError(f"Groq API not reachable — {e}")
    return "".join(parts).strip()


def _get_chat_model():
    global _chat_model
    if _chat_model is None and LANGGRAPH_OK:
        _chat_model = ChatGroq(
            model=GROQ_MODEL, api_key=GROQ_API_KEY, temperature=0.2,
        )
    return _chat_model


# ════════════════════════════════════════════════════════════
#  MODELS & SESSION
# ════════════════════════════════════════════════════════════

def _get_embed_model():
    global _embed_model
    if _embed_model is None:
        logger.info("Loading embedding model...")
        _embed_model = SentenceTransformer("all-mpnet-base-v2")
    return _embed_model


def _get_reranker():
    global _reranker
    if _reranker is None:
        logger.info("Loading reranker model...")
        _reranker = CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2")
    return _reranker


def _new_session(doc_id: str) -> dict:
    return {
        "doc_id": doc_id,
        "path": "",
        "title": "",
        "author": "",
        "page_count": 0,
        "segments": [],
        "chunks": [],
        "chunk_vectors": None,
        "tfidf": None,
        "tfidf_matrix": None,
        "document_text": "",
        "use_full": False,
        "history": [],
    }


def _active() -> dict | None:
    if _active_doc_id and _active_doc_id in _sessions:
        return _sessions[_active_doc_id]
    return None


# ════════════════════════════════════════════════════════════
#  HELPERS
# ════════════════════════════════════════════════════════════

def _extract_document_path(command: str) -> str | None:
    quoted = re.search(
        r'["\']([^"\']+\.(?:pdf|docx|doc|txt|text|md|pptx|ppt))["\']',
        command,
        re.IGNORECASE,
    )
    if quoted:
        return quoted.group(1)

    win_path = re.search(
        r'([A-Za-z]:\\(?:[^\\:*?"<>|\r\n]+\\)*[^\\:*?"<>|\r\n]+\.(?:pdf|docx|doc|txt|text|md|pptx|ppt))',
        command,
        re.IGNORECASE,
    )
    if win_path:
        return win_path.group(1)

    unix_path = re.search(
        r'(/(?:[^\s"\']+/)*[^\s"\']+\.(?:pdf|docx|doc|txt|text|md|pptx|ppt))',
        command,
        re.IGNORECASE,
    )
    if unix_path:
        return unix_path.group(1)

    rel_path = re.search(
        r'([^\s"\']+\.(?:pdf|docx|doc|txt|text|md|pptx|ppt))',
        command,
        re.IGNORECASE,
    )
    return rel_path.group(1) if rel_path else None


def _doc_id_from_path(path: str) -> str:
    return os.path.normcase(os.path.normpath(os.path.abspath(path)))


def _format_loc(loc: float) -> str:
    return f"Page {int(loc)}"


def _parse_loc(text: str) -> float | None:
    text = text.strip()
    m = re.match(r"^(?:page\s*)?(\d+)$", text, re.IGNORECASE)
    if m:
        return float(m.group(1))
    return None


def _fetch_document_metadata(path: str) -> dict:
    filename = os.path.basename(path)
    title = os.path.splitext(filename)[0]
    return {"title": title, "author": ""}


def _parse_pdf_segments(file_bytes: bytes) -> list:
    segments = []
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for i, page in enumerate(doc, 1):
            text = page.get_text().strip()
            if text:
                segments.append({"text": text, "start": float(i), "duration": 0})
    return segments


def _parse_docx_segments(file_bytes: bytes) -> list:
    doc = DocxDocument(io.BytesIO(file_bytes))
    segments = []
    idx = 1
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            segments.append({"text": text, "start": float(idx), "duration": 0})
            idx += 1
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                segments.append({"text": row_text, "start": float(idx), "duration": 0})
                idx += 1
    return segments


def _parse_pptx_segments(file_bytes: bytes) -> list:
    prs = PptxPresentation(io.BytesIO(file_bytes))
    segments = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(notes)
        if parts:
            segments.append({
                "text": "\n".join(parts),
                "start": float(i),
                "duration": 0,
            })
    return segments


def _parse_txt_segments(file_bytes: bytes) -> list:
    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = file_bytes.decode("latin-1", errors="replace")

    paragraphs = re.split(r"\n\s*\n", text.strip())
    segments = []
    for i, para in enumerate(paragraphs, 1):
        para = para.strip()
        if para:
            segments.append({"text": para, "start": float(i), "duration": 0})
    if not segments and text.strip():
        segments.append({"text": text.strip(), "start": 1.0, "duration": 0})
    return segments


def _get_document_segments(path: str):
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    try:
        with open(path, "rb") as f:
            file_bytes = f.read()
    except OSError as e:
        logger.error(f"Document read error: {e}")
        return None, "unreadable"

    try:
        if ext == "pdf":
            if not PDF_OK:
                return None, "pdf_deps"
            segments = _parse_pdf_segments(file_bytes)
        elif ext in ("docx", "doc"):
            if not DOCX_OK:
                return None, "docx_deps"
            segments = _parse_docx_segments(file_bytes)
        elif ext in ("pptx", "ppt"):
            if not PPTX_OK:
                return None, "pptx_deps"
            segments = _parse_pptx_segments(file_bytes)
        elif ext in ("txt", "text", "md"):
            segments = _parse_txt_segments(file_bytes)
        else:
            return None, "unsupported"
    except Exception as e:
        logger.error(f"Document parse error: {e}")
        return None, "unknown"

    if not segments:
        return None, "empty"
    return segments, None


def _sentence_aware_chunks_with_loc(segments: list, max_words: int = CHUNK_MAX_WORDS,
                                    overlap_sentences: int = CHUNK_OVERLAP_SENTS) -> list:
    words_with_loc = []
    for seg in segments:
        for w in seg["text"].split():
            words_with_loc.append((w, seg["start"]))

    full_text = " ".join(w for w, _ in words_with_loc)
    sentences = re.split(r"(?<=[.!?])\s+", full_text.strip())
    sentences = [s.strip() for s in sentences if s.strip()]

    chunks, word_idx = [], 0
    current_sentences, current_word_count = [], 0
    chunk_start_loc = words_with_loc[0][1] if words_with_loc else 1

    for sentence in sentences:
        sent_word_count = len(sentence.split())
        sent_start_loc = words_with_loc[word_idx][1] if word_idx < len(words_with_loc) else chunk_start_loc
        word_idx += sent_word_count

        if current_word_count + sent_word_count > max_words and current_sentences:
            chunk_end_loc = words_with_loc[min(word_idx - 1, len(words_with_loc) - 1)][1]
            chunks.append({
                "text": " ".join(current_sentences),
                "start": chunk_start_loc,
                "end": chunk_end_loc,
            })
            overlap = current_sentences[-overlap_sentences:]
            current_sentences = overlap
            current_word_count = sum(len(s.split()) for s in overlap)
            chunk_start_loc = sent_start_loc

        current_sentences.append(sentence)
        current_word_count += sent_word_count

    if current_sentences:
        chunk_end_loc = words_with_loc[-1][1] if words_with_loc else chunk_start_loc
        chunks.append({
            "text": " ".join(current_sentences),
            "start": chunk_start_loc,
            "end": chunk_end_loc,
        })
    return chunks


def _index_session(session: dict):
    chunks = session["chunks"]
    if len(chunks) <= 1:
        session["use_full"] = True
        return

    embed_model = _get_embed_model()
    texts = [c["text"] for c in chunks]
    session["chunk_vectors"] = embed_model.encode(texts, normalize_embeddings=True)
    session["tfidf"] = TfidfVectorizer(stop_words="english", max_features=50000)
    session["tfidf_matrix"] = session["tfidf"].fit_transform(texts)
    session["use_full"] = False


def _sample_chunks(chunks: list, n: int = 16) -> list:
    if len(chunks) <= n:
        return chunks
    step = max(1, len(chunks) // n)
    return chunks[::step][:n]


def _format_chunks(chunks: list, label: str = "Excerpt") -> str:
    return "\n\n".join(
        f"[{label} @ {_format_loc(c['start'])}–{_format_loc(c['end'])}]: {c['text']}"
        for c in chunks
    )


def _retrieve_chunks(session: dict, query: str, top_k: int = 12, final_k: int = 6) -> list:
    if session["use_full"]:
        return session["chunks"]

    chunks = session["chunks"]
    chunk_vectors = session["chunk_vectors"]
    tfidf = session["tfidf"]
    tfidf_matrix = session["tfidf_matrix"]

    embed_model = _get_embed_model()
    query_vec = embed_model.encode([query], normalize_embeddings=True)
    semantic_scores = np.dot(chunk_vectors, query_vec.T).flatten()

    query_tfidf = tfidf.transform([query])
    keyword_scores = cosine_similarity(query_tfidf, tfidf_matrix).flatten()

    hybrid = 0.6 * semantic_scores + 0.4 * keyword_scores
    top_indices = np.argsort(hybrid)[-top_k:][::-1]
    candidates = [(idx, chunks[idx]) for idx in top_indices]

    reranker = _get_reranker()
    pairs = [[query, c["text"]] for _, c in candidates]
    scores = reranker.predict(pairs)

    ranked = sorted(zip(scores, candidates), key=lambda x: x[0], reverse=True)
    result = [c for _, (_, c) in ranked[:final_k]]
    result.sort(key=lambda c: c["start"])
    return result


def _retrieve_context(session: dict, query: str, top_k: int = 12, final_k: int = 6) -> str:
    chunks = _retrieve_chunks(session, query, top_k, final_k)
    if session["use_full"] and len(chunks) == 1:
        c = chunks[0]
        return f"[Full document @ {_format_loc(c['start'])}–{_format_loc(c['end'])}]: {c['text']}"
    return _format_chunks(chunks)


def _wide_context(session: dict, query: str = "overview summary key points") -> str:
    chunks = _sample_chunks(session["chunks"], 20)
    if not session["use_full"]:
        extra = _retrieve_chunks(session, query, top_k=8, final_k=4)
        seen = {c["start"] for c in chunks}
        for c in extra:
            if c["start"] not in seen:
                chunks.append(c)
                seen.add(c["start"])
        chunks.sort(key=lambda c: c["start"])
    return _format_chunks(chunks)


# ════════════════════════════════════════════════════════════
#  DIRECT (NON-LLM) HANDLERS
# ════════════════════════════════════════════════════════════

def _segments_at_loc(session: dict, loc: float, window: float = 2) -> list:
    return [
        s for s in session["segments"]
        if loc - window <= s["start"] <= loc + window
    ]


def _content_at_page(session: dict, page_str: str) -> str:
    page = _parse_loc(page_str)
    if page is None:
        return "Invalid page. Use format: page 5 or just 5."
    segs = _segments_at_loc(session, page, window=0.5)
    if not segs:
        return NOT_IN_DOCUMENT
    lines = [f"[{_format_loc(s['start'])}] {s['text']}" for s in segs]
    return "\n".join(lines)


def _content_range(session: dict, start_str: str, end_str: str) -> str:
    start = _parse_loc(start_str)
    end = _parse_loc(end_str)
    if start is None or end is None:
        return "Invalid range. Use page numbers for start and end."
    if end < start:
        start, end = end, start
    segs = [s for s in session["segments"] if start <= s["start"] <= end]
    if not segs:
        return NOT_IN_DOCUMENT
    return "\n".join(f"[{_format_loc(s['start'])}] {s['text']}" for s in segs)


def _keyword_search(session: dict, keyword: str, max_hits: int = 15) -> str:
    keyword_lower = keyword.lower()
    hits = []
    for seg in session["segments"]:
        if keyword_lower in seg["text"].lower():
            hits.append(f"[{_format_loc(seg['start'])}] {seg['text']}")
            if len(hits) >= max_hits:
                break
    if not hits:
        return NOT_IN_DOCUMENT
    return f"Found '{keyword}' at:\n" + "\n".join(hits)


def _extract_urls(session: dict) -> str:
    pattern = r"https?://[^\s\]\)\"\'\,]+"
    found = []
    for seg in session["segments"]:
        for url in re.findall(pattern, seg["text"]):
            loc = _format_loc(seg["start"])
            entry = f"[{loc}] {url.rstrip('.,')}"
            if entry not in found:
                found.append(entry)
    if not found:
        return NOT_IN_DOCUMENT
    return "\n".join(found)


# ════════════════════════════════════════════════════════════
#  WEB (optional fact-check)
# ════════════════════════════════════════════════════════════

def _web_search_snippets(query: str, num: int = 3) -> str:
    if not WEB_SEARCH_OK:
        return "Web search unavailable. Install: pip install googlesearch-python"
    try:
        urls = list(google_search(query, num_results=num, advanced=False))
        return "\n".join(f"- {u}" for u in urls) if urls else "No web results found."
    except Exception as e:
        return f"Web search failed: {e}"


# ════════════════════════════════════════════════════════════
#  INTENT DETECTION
# ════════════════════════════════════════════════════════════

def _detect_intent(command: str) -> tuple[str, dict]:
    c = command.lower().strip()
    params: dict = {}

    # Compare two documents
    flat_paths = []
    for match in re.finditer(
        r'["\']([^"\']+\.(?:pdf|docx|doc|txt|text|md|pptx|ppt))["\']|'
        r'([A-Za-z]:\\[^\s"\']+\.(?:pdf|docx|doc|txt|text|md|pptx|ppt))|'
        r'(/(?:[^\s"\']+/)*[^\s"\']+\.(?:pdf|docx|doc|txt|text|md|pptx|ppt))|'
        r'([^\s"\']+\.(?:pdf|docx|doc|txt|text|md|pptx|ppt))',
        command,
        re.IGNORECASE,
    ):
        flat_paths.append(next(g for g in match.groups() if g))
    if len(flat_paths) >= 2 and any(w in c for w in ["compare", "vs", "versus", "difference"]):
        params["path_a"], params["path_b"] = flat_paths[0], flat_paths[1]
        return "compare_documents", params

    # Page range content
    range_m = re.search(
        r"(?:content|text|pages?)\s+(?:from|between)\s+(?:page\s*)?(\d+)\s+(?:to|and|-)\s+(?:page\s*)?(\d+)",
        c,
    )
    if range_m:
        params["start"], params["end"] = range_m.group(1), range_m.group(2)
        return "content_range", params

    # What is on a page
    at_m = re.search(
        r"(?:on|at|@|page)\s*(?:page\s*)?(\d+)|"
        r"what (?:is|was|does it say) (?:on|at) (?:page\s*)?(\d+)",
        c,
    )
    if at_m:
        params["page"] = at_m.group(1) or at_m.group(2)
        return "page_query", params

    # Keyword search
    kw_m = re.search(
        r"(?:search|find|where).*(?:mention|said|talk|discuss).*['\"]?(.+?)['\"]?$|"
        r"where is (.+?) mentioned",
        c,
    )
    if kw_m:
        params["keyword"] = (kw_m.group(1) or kw_m.group(2) or "").strip(" '\"?")
        if params["keyword"]:
            return "keyword_search", params

    # Translate
    tr_m = re.search(r"(?:translate|in)\s+(?:to\s+)?([a-zA-Z\u0900-\u097F]+)", c)
    if tr_m and any(w in c for w in ["translate", " hindi", " urdu", " spanish", " french", " language"]):
        params["lang"] = tr_m.group(1)
        return "translate", params

    # Fact-check / up-to-date (needs web)
    if any(w in c for w in ["fact check", "fact-check", "verify", "is this true", "is it true"]):
        return "fact_check", params
    if any(w in c for w in ["up to date", "uptodate", "outdated", "still valid", "current"]):
        return "freshness", params

    intent_map = [
        ("summarize", ["summarize", "summary of the document", "give me a summary"]),
        ("main_topic", ["main topic", "what is this document about", "what is the topic"]),
        ("explain_simple", ["explain simply", "simple language", "eli5", "easy words"]),
        ("key_points", ["key points", "main points", "important points"]),
        ("takeaways", ["takeaways", "take away", "lessons learned"]),
        ("tips", ["tips", "tricks", "strategies", "hacks"]),
        ("resources", ["tools", "websites", "products", "books", "papers", "resources mentioned"]),
        ("examples", ["with examples", "give examples", "example"]),
        ("conclusion", ["conclusion", "wrap up", "final thoughts"]),
        ("short_notes", ["short notes", "brief notes"]),
        ("bullet_summary", ["bullet", "bullet point", "point wise"]),
        ("blog_post", ["blog post", "write a blog", "convert to blog"]),
        ("study_notes", ["study notes", "revision notes"]),
        ("interview", ["interview questions", "interview q"]),
        ("quiz", ["quiz", "mcq", "multiple choice"]),
        ("facts", ["facts", "numbers", "statistics", "stats", "data mentioned"]),
        ("contradictions", ["contradiction", "inconsistent", "conflict"]),
        ("beginner", ["for beginners", "beginner friendly", "i am a beginner"]),
        ("expert", ["for experts", "advanced", "expert level"]),
        ("action_items", ["action items", "action plan", "what should i do", "next steps"]),
        ("code", ["code snippet", "commands", "api", "syntax", "terminal command"]),
        ("links", ["links", "urls", "references", "external link"]),
        ("questions_answered", ["what questions does this document answer", "questions answered"]),
        ("questions_not_answered", ["what questions does this document not answer", "not answer", "left unanswered"]),
        ("entities", ["people mentioned", "companies mentioned", "products mentioned", "technologies mentioned"]),
        ("workflow", ["workflow", "step by step", "steps", "process", "procedure"]),
        ("checklist", ["checklist", "check list"]),
        ("flashcards", ["flashcards", "flash cards"]),
        ("sections", ["each section", "section by section", "break down sections"]),
        ("insights", ["insights", "missed", "what did i miss", "important insights"]),
        ("pages", ["page for", "which page", "on what page", "page number"]),
    ]

    for intent, keywords in intent_map:
        if any(k in c for k in keywords):
            return intent, params

    return "qa", params


# ════════════════════════════════════════════════════════════
#  PROMPTS
# ════════════════════════════════════════════════════════════

_SYSTEM_GROUNDED = """You are a document assistant. Use ONLY the document excerpts provided.
Rules:
- Cite page numbers like [Page N] whenever possible.
- Be clear and direct. No filler or disclaimers.
- If information is missing, reply exactly: "This information is not available in the document."
- Do not invent facts, links, names, or numbers."""

_INTENT_INSTRUCTIONS = {
    "summarize": "Summarize the document in 3-5 sentences plus bullet key points with page numbers.",
    "main_topic": "State the main topic in 1-2 sentences with the primary page reference.",
    "explain_simple": "Explain the main concepts in simple language with page numbers.",
    "key_points": "List key points as bullets. Each bullet: point + [Page N].",
    "takeaways": "List important takeaways as bullets with page numbers.",
    "tips": "Extract all tips, tricks, and strategies with page numbers.",
    "resources": "List every tool, website, product, book, paper, or resource mentioned with page numbers.",
    "examples": "Explain concepts using examples from the document with page numbers.",
    "conclusion": "Write a concise conclusion of the document with page numbers.",
    "short_notes": "Create short revision notes with page numbers.",
    "bullet_summary": "Bullet-point summary covering the full document with page numbers.",
    "blog_post": "Convert the document into a structured blog post (title, intro, sections, conclusion) using only document content.",
    "study_notes": "Create study notes: headings, definitions, key facts, each with page numbers.",
    "interview": "Generate 10 interview questions (+ brief expected answers from the document) with page numbers.",
    "quiz": "Generate 5 MCQs with 4 options each, mark correct answer, cite page for each.",
    "facts": "Extract all facts, numbers, and statistics mentioned with page numbers.",
    "contradictions": "Identify any contradictions or inconsistencies in the document with page numbers. If none, say so.",
    "beginner": "Explain the document content for a complete beginner with page numbers.",
    "expert": "Explain the document content for an expert audience with page numbers.",
    "action_items": "Generate actionable steps the reader should take based on the document.",
    "code": "Extract all code snippets, commands, and APIs mentioned with page numbers.",
    "links": "List all external links and references mentioned (URLs if written, otherwise describe).",
    "questions_answered": "List questions this document answers (as bullet questions) with page numbers.",
    "questions_not_answered": "List important questions the document does NOT answer.",
    "entities": "List people, companies, products, and technologies mentioned with page numbers.",
    "workflow": "Extract the workflow or step-by-step process with numbered steps and page numbers.",
    "checklist": "Convert content into a practical checklist with page numbers.",
    "flashcards": "Generate 10 flashcards (Q on one line, A on next) with page numbers.",
    "sections": "Explain each major section of the document one by one with page numbers.",
    "insights": "Suggest important insights the reader may have missed, with page numbers.",
    "pages": "For the topic in the user question, list which pages discuss it with page numbers and brief quotes.",
    "qa": "Answer the question directly. Shortest accurate answer first, then page citation if relevant.",
    "translate": "Answer in the requested language. Ground in document only. Include page numbers.",
    "compare_documents": "Compare the two documents: topic, approach, key differences, overlap. Cite pages from each.",
    "fact_check": "Fact-check claims from the document using web sources provided. Note what the document says vs external info.",
    "freshness": "Assess whether the document's information appears up to date based on content and web hints.",
}


def _build_prompt(intent: str, context: str, question: str, extra: str = "") -> str:
    instruction = _INTENT_INSTRUCTIONS.get(intent, _INTENT_INSTRUCTIONS["qa"])
    meta = f"\n{extra}\n" if extra else ""
    return f"""Task: {instruction}

User request: {question}
{meta}
--- DOCUMENT EXCERPTS ---
{context}
----------------------------

Response:"""


# ════════════════════════════════════════════════════════════
#  CORE ASK ENGINE
# ════════════════════════════════════════════════════════════

def _run_task(session: dict, command: str, intent: str, params: dict,
              stream: bool = False, on_token: Callable[[str], None] | None = None,
              history_snippet: str = "") -> str:
    # Direct handlers (no LLM)
    if intent == "page_query":
        return _content_at_page(session, params["page"])
    if intent == "content_range":
        return _content_range(session, params["start"], params["end"])
    if intent == "keyword_search":
        return _keyword_search(session, params["keyword"])
    if intent == "links":
        urls = _extract_urls(session)
        if urls != NOT_IN_DOCUMENT:
            return urls

    if not is_ollama_ready():
        return _ollama_error()

    # Context selection by intent
    wide_intents = {
        "summarize", "main_topic", "bullet_summary", "blog_post", "study_notes",
        "sections", "conclusion", "questions_answered", "questions_not_answered",
        "contradictions", "insights", "compare_documents", "checklist", "workflow",
    }
    if intent in wide_intents:
        context = _wide_context(session, command)
    elif intent in ("facts", "resources", "code", "entities", "tips"):
        context = _retrieve_context(session, command, top_k=16, final_k=10)
    else:
        context = _retrieve_context(session, command)

    extra = ""
    if session.get("title"):
        extra += f"Document title: {session['title']}\n"
    if session.get("path"):
        extra += f"File: {session['path']}\n"
    if session.get("page_count"):
        extra += f"Pages/sections: {session['page_count']}\n"
    if history_snippet:
        extra += f"\n{history_snippet}\n"

    if intent == "compare_documents":
        path_b = params.get("path_b", "")
        doc_b = _doc_id_from_path(os.path.abspath(os.path.expanduser(path_b))) if path_b else None
        if doc_b and doc_b in _sessions:
            ctx_b = _wide_context(_sessions[doc_b], command)
            extra += f"\n--- SECOND DOCUMENT EXCERPTS ---\n{ctx_b}\n"

    if intent in ("fact_check", "freshness"):
        claim = _retrieve_context(session, command, top_k=6, final_k=4)
        web = _web_search_snippets(command[:120])
        extra += f"\nWeb search results (for verification only):\n{web}\n"
        context = claim

    prompt = _build_prompt(intent, context, command, extra)

    try:
        if stream and on_token:
            return ask_ollama_stream(prompt, temperature=0.2, system=_SYSTEM_GROUNDED, on_token=on_token)
        return ask_ollama(prompt, temperature=0.2, system=_SYSTEM_GROUNDED)
    except RuntimeError as e:
        logger.error(str(e))
        return str(e)


def process_document_query(command: str, doc_id: str | None = None,
                           stream: bool = False,
                           on_token: Callable[[str], None] | None = None) -> str:
    session = _sessions.get(doc_id) if doc_id else _active()
    if not session:
        return "No document loaded. Provide a document path first."

    intent, params = _detect_intent(command)

    # Compare: auto-load second document if path given
    if intent == "compare_documents" and params.get("path_b"):
        path_b = params["path_b"]
        doc_b = _doc_id_from_path(os.path.abspath(os.path.expanduser(path_b)))
        if doc_b not in _sessions:
            load_document(f"load {path_b}")

    answer = _run_task(session, command, intent, params, stream=stream, on_token=on_token)
    session["history"].append((command, answer))
    return answer


# ════════════════════════════════════════════════════════════
#  LANGGRAPH MEMORY (optional)
# ════════════════════════════════════════════════════════════

if LANGGRAPH_OK:
    class DocumentChat(TypedDict):
        messages: Annotated[list, add_messages]
        doc_id: str


def _make_document_graph(doc_id: str):
    def response_generate(state):
        user_query = state["messages"][-1].content
        answer = process_document_query(user_query, doc_id=state.get("doc_id", doc_id))
        return {"messages": [HumanMessage(content=answer)]}

    graph = StateGraph(DocumentChat)
    graph.add_node("response_generate", response_generate)
    graph.add_edge(START, "response_generate")
    graph.add_edge("response_generate", END)
    return graph.compile(checkpointer=MemorySaver())


def _get_or_create_document_graph(doc_id: str):
    if doc_id not in _graph_store:
        _graph_store[doc_id] = {
            "graph": _make_document_graph(doc_id),
            "config": {"configurable": {"thread_id": doc_id}},
        }
    return _graph_store[doc_id]


def ask_document_langgraph(doc_id: str, question: str,
                           stream: bool = False,
                           on_token: Callable[[str], None] | None = None) -> str:
    if not LANGGRAPH_OK:
        return process_document_query(question, doc_id=doc_id, stream=stream, on_token=on_token)
    if doc_id not in _sessions:
        return "No document loaded yet."
    if stream:
        return process_document_query(question, doc_id=doc_id, stream=True, on_token=on_token)

    bot = _get_or_create_document_graph(doc_id)
    result = bot["graph"].invoke(
        {"messages": [HumanMessage(content=question)], "doc_id": doc_id},
        bot["config"],
    )
    return result["messages"][-1].content


# ════════════════════════════════════════════════════════════
#  PUBLIC API
# ════════════════════════════════════════════════════════════

def load_document(command: str) -> str:
    global _active_doc_id

    if not DOC_ML_OK:
        msg = "Missing deps: pip install sentence-transformers scikit-learn numpy"
        speak(msg)
        return msg

    path = _extract_document_path(command)
    if not path:
        msg = "Please provide the document file path."
        speak(msg)
        return msg

    path = os.path.normpath(os.path.abspath(os.path.expanduser(path)))
    if not os.path.isfile(path):
        msg = f"File not found: {path}"
        speak(msg)
        return msg

    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf" and not PDF_OK:
        msg = "PDF support missing. Run: pip install pymupdf"
        speak(msg)
        return msg
    if ext in (".docx", ".doc") and not DOCX_OK:
        msg = "DOCX support missing. Run: pip install python-docx"
        speak(msg)
        return msg
    if ext in (".pptx", ".ppt") and not PPTX_OK:
        msg = "PPTX support missing. Run: pip install python-pptx"
        speak(msg)
        return msg

    doc_id = _doc_id_from_path(path)
    if doc_id in _sessions:
        _active_doc_id = doc_id
        msg = f"Document already loaded — switched to {os.path.basename(path)}."
        speak(msg)
        return msg

    segments, err = _get_document_segments(path)
    if segments is None:
        reasons = {
            "unreadable": f"Could not read file: {path}",
            "pdf_deps": "PDF support missing. Run: pip install pymupdf",
            "docx_deps": "DOCX support missing. Run: pip install python-docx",
            "pptx_deps": "PPTX support missing. Run: pip install python-pptx",
            "unsupported": "Unsupported file type. Use PDF, DOCX, TXT, or PPTX.",
            "empty": "No text could be extracted from the document.",
            "unknown": "Could not parse document. Try again.",
        }
        msg = reasons.get(err, "Document load failed.")
        speak(msg)
        return msg

    meta = _fetch_document_metadata(path)
    document_text = " ".join(s["text"] for s in segments)
    word_count = len(document_text.split())
    page_count = len(segments)

    session = _new_session(doc_id)
    session["path"] = path
    session["segments"] = segments
    session["document_text"] = document_text
    session["title"] = meta["title"]
    session["author"] = meta["author"]
    session["page_count"] = page_count

    if word_count <= FULL_DOCUMENT_WORD_LIMIT:
        session["chunks"] = [{
            "text": document_text,
            "start": segments[0]["start"] if segments else 1,
            "end": segments[-1]["start"] if segments else 1,
        }]
        session["use_full"] = True
    else:
        session["chunks"] = _sentence_aware_chunks_with_loc(segments)
        _index_session(session)

    _sessions[doc_id] = session
    _active_doc_id = doc_id

    title_bit = f' "{meta["title"]}"' if meta["title"] else ""
    msg = (
        f"Loaded{title_bit} — {page_count} pages/sections, {word_count} words, "
        f"{len(session['chunks'])} chunks indexed. "
        f"Ask anything: summarize, key points, quiz, pages, facts, etc."
    )
    logger.info(f"Document loaded: {path} ({word_count} words, {len(session['chunks'])} chunks)")
    return msg


def summarize_document(command: str = "") -> str:
    return process_document_query(command or "summarize the document")


def ask_document(command: str) -> str:
    user_query = re.sub(
        r"^(document|doc|file|pdf)\s+(question|ask|kya|mein)\s*", "",
        command, flags=re.IGNORECASE,
    ).strip() or command.strip()
    return process_document_query(user_query)


def list_loaded_documents() -> str:
    if not _sessions:
        return "No documents loaded."
    lines = []
    for did, s in _sessions.items():
        marker = " (active)" if did == _active_doc_id else ""
        title = s.get("title") or os.path.basename(s.get("path", did))
        lines.append(f"- {title}{marker} [{s.get('path', did)}] — {len(s['chunks'])} chunks")
    return "Loaded documents:\n" + "\n".join(lines)


def switch_document(command: str) -> str:
    global _active_doc_id
    path = _extract_document_path(command)
    candidate = path or command.strip()
    doc_id = _doc_id_from_path(os.path.abspath(os.path.expanduser(candidate))) if path else candidate

    if doc_id not in _sessions:
        return "Document not loaded. Use: load document <path>"
    _active_doc_id = doc_id
    return f"Switched to {os.path.basename(_sessions[doc_id].get('path', doc_id))}."


def document_chatbot(command: str) -> str:
    c = command.lower()

    wants_load = any(w in c for w in ["load", "kholo", "open"]) and \
                 ("document" in c or "doc" in c or "file" in c or "pdf" in c)
    if wants_load:
        path = _extract_document_path(command)
        if path:
            return load_document(command)
        return "Please provide the document file path."

    if any(w in c for w in ["switch", "badlo"]) and ("document" in c or "doc" in c or "file" in c):
        return switch_document(command)

    if any(w in c for w in ["list document", "loaded document"]):
        return list_loaded_documents()

    if not _active():
        path = _extract_document_path(command)
        if path:
            return load_document(command)
        return "Load a document first: load document <path>"

    return ask_document(command)


# ════════════════════════════════════════════════════════════════════════
#  WEB FEATURE — Flask Blueprint + SQLite persistence + LangGraph Threads
#  Powers templates/document_chatbot.html, static/document_chatbot.css and
#  js/document_chatbot.js. Reuses the extraction / chunking / retrieval /
#  prompting engine defined above, but keyed by a browser-issued Thread ID
#  instead of a filesystem path (one document + one conversation per
#  LangGraph Thread, exactly like a New Chat in ChatGPT).
# ════════════════════════════════════════════════════════════════════════

# ────────────────────────────────────────────────────────────────────────
#  SQLITE — threads + messages
#  (own database file: db_storage/document_chatbot.db — does not touch
#   db.py / snetch.db used by the rest of the app)
# ────────────────────────────────────────────────────────────────────────

def _doc_conn() -> sqlite3.Connection:
    os.makedirs(DB_STORAGE_DIR, exist_ok=True)
    conn = sqlite3.connect(DOC_DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def init_document_chatbot_db() -> None:
    os.makedirs(DB_STORAGE_DIR, exist_ok=True)
    os.makedirs(UPLOAD_DIR, exist_ok=True)
    conn = _doc_conn()
    conn.execute("""
        CREATE TABLE IF NOT EXISTS threads (
            thread_id         TEXT PRIMARY KEY,
            title             TEXT NOT NULL DEFAULT 'New Chat',
            doc_name          TEXT,
            doc_type          TEXT,
            doc_size          INTEGER,
            doc_path          TEXT,
            upload_time       INTEGER,
            upload_status     TEXT NOT NULL DEFAULT 'pending',
            processing_status TEXT NOT NULL DEFAULT 'awaiting_upload',
            is_pinned         INTEGER NOT NULL DEFAULT 0,
            is_archived       INTEGER NOT NULL DEFAULT 0,
            created_at        INTEGER NOT NULL,
            updated_at        INTEGER NOT NULL
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS messages (
            id         INTEGER PRIMARY KEY AUTOINCREMENT,
            thread_id  TEXT NOT NULL,
            role       TEXT NOT NULL,
            content    TEXT NOT NULL,
            liked      INTEGER NOT NULL DEFAULT 0,
            disliked   INTEGER NOT NULL DEFAULT 0,
            created_at INTEGER NOT NULL
        )
    """)
    conn.execute("CREATE INDEX IF NOT EXISTS idx_messages_thread ON messages(thread_id)")
    conn.commit()
    conn.close()


def create_thread_row(thread_id: str, title: str = "New Chat") -> None:
    now = int(time.time())
    conn = _doc_conn()
    conn.execute(
        "INSERT INTO threads (thread_id, title, created_at, updated_at) VALUES (?,?,?,?)",
        (thread_id, title, now, now),
    )
    conn.commit()
    conn.close()


def get_thread_row(thread_id: str):
    conn = _doc_conn()
    row = conn.execute("SELECT * FROM threads WHERE thread_id=?", (thread_id,)).fetchone()
    conn.close()
    return row


def list_thread_rows(search: str = "", archived: bool = False):
    conn = _doc_conn()
    if search:
        like = f"%{search.lower()}%"
        rows = conn.execute("""
            SELECT DISTINCT t.* FROM threads t
            LEFT JOIN messages m ON m.thread_id = t.thread_id
            WHERE t.is_archived = ?
              AND (LOWER(t.title) LIKE ? OR LOWER(COALESCE(t.doc_name,'')) LIKE ?
                   OR LOWER(m.content) LIKE ?)
            ORDER BY t.is_pinned DESC, t.updated_at DESC
        """, (1 if archived else 0, like, like, like)).fetchall()
    else:
        rows = conn.execute("""
            SELECT * FROM threads WHERE is_archived = ?
            ORDER BY is_pinned DESC, updated_at DESC
        """, (1 if archived else 0,)).fetchall()
    conn.close()
    return rows


def touch_thread(thread_id: str) -> None:
    conn = _doc_conn()
    conn.execute("UPDATE threads SET updated_at=? WHERE thread_id=?", (int(time.time()), thread_id))
    conn.commit()
    conn.close()


def set_thread_doc_info(thread_id: str, name: str, ftype: str, size: int, path: str,
                         upload_time: int, processing_status: str, upload_status: str) -> None:
    conn = _doc_conn()
    conn.execute("""
        UPDATE threads SET doc_name=?, doc_type=?, doc_size=?, doc_path=?, upload_time=?,
               processing_status=?, upload_status=?, updated_at=? WHERE thread_id=?
    """, (name, ftype, size, path, upload_time, processing_status, upload_status,
          int(time.time()), thread_id))
    conn.commit()
    conn.close()


def set_processing_status(thread_id: str, status: str) -> None:
    conn = _doc_conn()
    conn.execute("UPDATE threads SET processing_status=?, updated_at=? WHERE thread_id=?",
                 (status, int(time.time()), thread_id))
    conn.commit()
    conn.close()


def rename_thread_row(thread_id: str, title: str) -> None:
    conn = _doc_conn()
    conn.execute("UPDATE threads SET title=?, updated_at=? WHERE thread_id=?",
                 (title.strip()[:120] or "New Chat", int(time.time()), thread_id))
    conn.commit()
    conn.close()


def set_pin_row(thread_id: str, pinned: bool) -> None:
    conn = _doc_conn()
    conn.execute("UPDATE threads SET is_pinned=?, updated_at=? WHERE thread_id=?",
                 (1 if pinned else 0, int(time.time()), thread_id))
    conn.commit()
    conn.close()


def set_archive_row(thread_id: str, archived: bool) -> None:
    conn = _doc_conn()
    conn.execute("UPDATE threads SET is_archived=?, updated_at=? WHERE thread_id=?",
                 (1 if archived else 0, int(time.time()), thread_id))
    conn.commit()
    conn.close()


def delete_thread_row(thread_id: str) -> None:
    conn = _doc_conn()
    conn.execute("DELETE FROM threads WHERE thread_id=?", (thread_id,))
    conn.execute("DELETE FROM messages WHERE thread_id=?", (thread_id,))
    conn.commit()
    conn.close()


def add_message_row(thread_id: str, role: str, content: str) -> int:
    conn = _doc_conn()
    cur = conn.execute(
        "INSERT INTO messages (thread_id, role, content, created_at) VALUES (?,?,?,?)",
        (thread_id, role, content, int(time.time())),
    )
    conn.commit()
    msg_id = cur.lastrowid
    conn.close()
    return msg_id


def get_message_rows(thread_id: str, limit: int | None = None):
    conn = _doc_conn()
    if limit:
        rows = conn.execute(
            "SELECT * FROM messages WHERE thread_id=? ORDER BY id DESC LIMIT ?",
            (thread_id, limit),
        ).fetchall()
        rows = list(reversed(rows))
    else:
        rows = conn.execute(
            "SELECT * FROM messages WHERE thread_id=? ORDER BY id ASC", (thread_id,)
        ).fetchall()
    conn.close()
    return rows


def delete_message_row(message_id: int) -> None:
    conn = _doc_conn()
    conn.execute("DELETE FROM messages WHERE id=?", (message_id,))
    conn.commit()
    conn.close()


def set_message_feedback_row(message_id: int, feedback: str) -> None:
    liked = 1 if feedback == "like" else 0
    disliked = 1 if feedback == "dislike" else 0
    conn = _doc_conn()
    conn.execute("UPDATE messages SET liked=?, disliked=? WHERE id=?",
                 (liked, disliked, message_id))
    conn.commit()
    conn.close()


def _serialize_thread(row) -> dict:
    return {
        "thread_id": row["thread_id"],
        "title": row["title"],
        "document": {
            "name": row["doc_name"],
            "type": row["doc_type"],
            "size": row["doc_size"],
            "upload_time": row["upload_time"],
        } if row["doc_name"] else None,
        "upload_status": row["upload_status"],
        "processing_status": row["processing_status"],
        "is_pinned": bool(row["is_pinned"]),
        "is_archived": bool(row["is_archived"]),
        "created_at": row["created_at"],
        "updated_at": row["updated_at"],
    }


def _serialize_message(row) -> dict:
    return {
        "id": row["id"],
        "role": row["role"],
        "content": row["content"],
        "liked": bool(row["liked"]),
        "disliked": bool(row["disliked"]),
        "created_at": row["created_at"],
    }


# ────────────────────────────────────────────────────────────────────────
#  PER-THREAD DOCUMENT SESSION + RETRIEVAL (reuses the engine above)
# ────────────────────────────────────────────────────────────────────────

def _build_history_snippet(thread_id: str, limit: int = 6) -> str:
    rows = get_message_rows(thread_id, limit=limit)
    if not rows:
        return ""
    lines = []
    for r in rows:
        who = "User" if r["role"] == "user" else "Assistant"
        lines.append(f"{who}: {r['content'][:600]}")
    return "Previous conversation turns (context only — answer the NEW question below):\n" + "\n".join(lines)


def _load_web_document(thread_id: str, path: str, orig_name: str):
    """Parse + (optionally) index an uploaded document for a chat thread.
    Returns (ok: bool, info_or_error)."""
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    if ext == "pdf" and not PDF_OK:
        return False, "PDF support missing on the server. Install: pip install pymupdf"
    if ext in ("docx", "doc") and not DOCX_OK:
        return False, "DOCX support missing on the server. Install: pip install python-docx"
    if ext in ("pptx", "ppt") and not PPTX_OK:
        return False, "PPTX support missing on the server. Install: pip install python-pptx"

    segments, err = _get_document_segments(path)
    if segments is None:
        reasons = {
            "unreadable": "Could not read the uploaded file.",
            "pdf_deps": "PDF support missing on the server.",
            "docx_deps": "DOCX support missing on the server.",
            "pptx_deps": "PPTX support missing on the server.",
            "unsupported": "Unsupported file type. Use PDF, DOCX, TXT, Markdown or PPTX.",
            "empty": "No text could be extracted from the document.",
            "unknown": "Could not parse the document. Please try again.",
        }
        return False, reasons.get(err, "Document processing failed.")

    meta = _fetch_document_metadata(path)
    document_text = " ".join(s["text"] for s in segments)
    word_count = len(document_text.split())

    session = _new_session(thread_id)
    session["path"] = path
    session["segments"] = segments
    session["document_text"] = document_text
    session["title"] = meta["title"] or orig_name
    session["page_count"] = len(segments)

    if word_count <= FULL_DOCUMENT_WORD_LIMIT:
        session["chunks"] = [{
            "text": document_text,
            "start": segments[0]["start"] if segments else 1,
            "end": segments[-1]["start"] if segments else 1,
        }]
        session["use_full"] = True
    else:
        if not DOC_ML_OK:
            return False, ("This document is large and needs semantic indexing, but the "
                            "server is missing: sentence-transformers, scikit-learn, numpy.")
        session["chunks"] = _sentence_aware_chunks_with_loc(segments)
        _index_session(session)

    _web_sessions[thread_id] = session
    return True, {"word_count": word_count, "page_count": len(segments), "chunks": len(session["chunks"])}


def process_web_query(thread_id: str, question: str, stream: bool = False,
                       on_token: Callable[[str], None] | None = None) -> str:
    session = _web_sessions.get(thread_id)
    if not session:
        raise ValueError("No document is loaded for this chat yet.")
    intent, params = _detect_intent(question)
    history_snippet = _build_history_snippet(thread_id)
    return _run_task(session, question, intent, params, stream=stream,
                      on_token=on_token, history_snippet=history_snippet)


# ────────────────────────────────────────────────────────────────────────
#  LANGGRAPH THREAD MEMORY — synced with SQLite via SqliteSaver when the
#  optional langgraph / langchain packages are installed. Every New Chat
#  gets its own Thread ID and its own isolated checkpointed memory; no
#  Thread ever inherits memory from another. Falls back gracefully (the
#  feature still works end-to-end, using the messages table above as the
#  source of truth) when those optional heavy packages are not installed.
# ────────────────────────────────────────────────────────────────────────

if LANGGRAPH_OK:
    class WebDocChat(TypedDict):
        messages: Annotated[list, add_messages]
        thread_id: str

    def _get_checkpointer():
        global _checkpointer_singleton
        if _checkpointer_singleton is not None:
            return _checkpointer_singleton
        try:
            from langgraph.checkpoint.sqlite import SqliteSaver
            conn = sqlite3.connect(CHECKPOINT_DB_PATH, check_same_thread=False)
            _checkpointer_singleton = SqliteSaver(conn)
        except ImportError:
            _checkpointer_singleton = MemorySaver()
        return _checkpointer_singleton

    def _web_graph_node(state):
        thread_id = state["thread_id"]
        question = state["messages"][-1].content
        answer = process_web_query(thread_id, question, stream=False)
        return {"messages": [HumanMessage(content=answer)]}

    def _make_web_graph():
        g = StateGraph(WebDocChat)
        g.add_node("answer", _web_graph_node)
        g.add_edge(START, "answer")
        g.add_edge("answer", END)
        return g.compile(checkpointer=_get_checkpointer())

    def _get_or_create_web_graph(thread_id: str):
        if thread_id not in _web_graph_store:
            _web_graph_store[thread_id] = {
                "graph": _make_web_graph(),
                "config": {"configurable": {"thread_id": thread_id}},
            }
        return _web_graph_store[thread_id]

    def _sync_langgraph_memory(thread_id: str, question: str, answer: str) -> None:
        """Persist this Q/A turn into the Thread's LangGraph checkpoint
        (SQLite-backed) without re-running the LLM — the answer was
        already streamed to the user."""
        try:
            bot = _get_or_create_web_graph(thread_id)
            bot["graph"].update_state(
                bot["config"],
                {"messages": [HumanMessage(content=question), HumanMessage(content=answer)]},
            )
        except Exception as e:
            logger.warning(f"LangGraph memory sync skipped: {e}")

    def _drop_langgraph_thread(thread_id: str) -> None:
        _web_graph_store.pop(thread_id, None)
else:
    def _sync_langgraph_memory(thread_id: str, question: str, answer: str) -> None:
        pass

    def _drop_langgraph_thread(thread_id: str) -> None:
        pass


# ────────────────────────────────────────────────────────────────────────
#  FLASK BLUEPRINT — /document_chatbot/api/*
# ────────────────────────────────────────────────────────────────────────

document_chatbot_bp = Blueprint(
    "document_chatbot_api", __name__, url_prefix="/document_chatbot/api"
)


@document_chatbot_bp.route("/new_chat", methods=["POST"])
def api_new_chat():
    thread_id = uuid.uuid4().hex
    create_thread_row(thread_id, "New Chat")
    return jsonify({"success": True, "thread_id": thread_id,
                     "thread": _serialize_thread(get_thread_row(thread_id))})


@document_chatbot_bp.route("/upload", methods=["POST"])
def api_upload():
    thread_id = request.form.get("thread_id", "").strip()
    file = request.files.get("file")

    if not thread_id or not get_thread_row(thread_id):
        return jsonify({"success": False, "error": "Invalid or missing chat thread."}), 400
    if not file or not file.filename:
        return jsonify({"success": False, "error": "No file provided."}), 400

    filename = secure_filename(file.filename) or "document"
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in ALLOWED_EXTENSIONS:
        return jsonify({"success": False, "error": "Unsupported file type. Use PDF, DOCX, TXT, "
                                                     "Markdown or PPTX."}), 400

    set_processing_status(thread_id, "uploading")
    os.makedirs(UPLOAD_DIR, exist_ok=True)
    dest_path = os.path.join(UPLOAD_DIR, f"{thread_id}.{ext}")

    file.save(dest_path)
    size = os.path.getsize(dest_path)
    if size > MAX_UPLOAD_BYTES:
        os.remove(dest_path)
        set_processing_status(thread_id, "failed")
        return jsonify({"success": False, "error": "File exceeds the 50MB limit."}), 400
    if size == 0:
        os.remove(dest_path)
        set_processing_status(thread_id, "failed")
        return jsonify({"success": False, "error": "Uploaded file is empty."}), 400

    upload_time = int(time.time())
    set_thread_doc_info(thread_id, filename, ext, size, dest_path, upload_time,
                         "processing", "success")

    ok, info = _load_web_document(thread_id, dest_path, filename)
    if not ok:
        set_processing_status(thread_id, "failed")
        return jsonify({"success": False, "error": info,
                         "processing_status": "failed"}), 422

    set_processing_status(thread_id, "ready")

    row = get_thread_row(thread_id)
    if row["title"] == "New Chat":
        rename_thread_row(thread_id, os.path.splitext(filename)[0][:80])

    touch_thread(thread_id)
    return jsonify({
        "success": True,
        "thread": _serialize_thread(get_thread_row(thread_id)),
        "stats": info,
    })


@document_chatbot_bp.route("/threads", methods=["GET"])
def api_threads():
    search = request.args.get("search", "").strip()
    scope = request.args.get("scope", "all")
    rows = list_thread_rows(search=search, archived=(scope == "archived"))
    threads = [_serialize_thread(r) for r in rows]
    pinned = [t for t in threads if t["is_pinned"]]
    recent = [t for t in threads if not t["is_pinned"]]
    return jsonify({"success": True, "pinned": pinned, "recent": recent})


@document_chatbot_bp.route("/thread/<thread_id>", methods=["GET"])
def api_thread_detail(thread_id):
    row = get_thread_row(thread_id)
    if not row:
        return jsonify({"success": False, "error": "Chat not found."}), 404
    msgs = get_message_rows(thread_id)
    return jsonify({
        "success": True,
        "thread": _serialize_thread(row),
        "messages": [_serialize_message(m) for m in msgs],
    })


@document_chatbot_bp.route("/thread/<thread_id>/rename", methods=["POST"])
def api_rename_thread(thread_id):
    if not get_thread_row(thread_id):
        return jsonify({"success": False, "error": "Chat not found."}), 404
    data = request.get_json(force=True, silent=True) or {}
    title = (data.get("title") or "").strip()
    if not title:
        return jsonify({"success": False, "error": "Title cannot be empty."}), 400
    rename_thread_row(thread_id, title)
    return jsonify({"success": True, "thread": _serialize_thread(get_thread_row(thread_id))})


@document_chatbot_bp.route("/thread/<thread_id>/pin", methods=["POST"])
def api_pin_thread(thread_id):
    row = get_thread_row(thread_id)
    if not row:
        return jsonify({"success": False, "error": "Chat not found."}), 404
    data = request.get_json(force=True, silent=True) or {}
    pinned = data.get("pinned")
    pinned = (not row["is_pinned"]) if pinned is None else bool(pinned)
    set_pin_row(thread_id, pinned)
    return jsonify({"success": True, "thread": _serialize_thread(get_thread_row(thread_id))})


@document_chatbot_bp.route("/thread/<thread_id>/archive", methods=["POST"])
def api_archive_thread(thread_id):
    row = get_thread_row(thread_id)
    if not row:
        return jsonify({"success": False, "error": "Chat not found."}), 404
    data = request.get_json(force=True, silent=True) or {}
    archived = data.get("archived")
    archived = (not row["is_archived"]) if archived is None else bool(archived)
    set_archive_row(thread_id, archived)
    return jsonify({"success": True, "thread": _serialize_thread(get_thread_row(thread_id))})


@document_chatbot_bp.route("/thread/<thread_id>", methods=["DELETE"])
def api_delete_thread(thread_id):
    row = get_thread_row(thread_id)
    if not row:
        return jsonify({"success": False, "error": "Chat not found."}), 404
    if row["doc_path"] and os.path.isfile(row["doc_path"]):
        try:
            os.remove(row["doc_path"])
        except OSError:
            pass
    delete_thread_row(thread_id)
    _web_sessions.pop(thread_id, None)
    _drop_langgraph_thread(thread_id)
    return jsonify({"success": True})


@document_chatbot_bp.route("/thread/<thread_id>/download", methods=["GET"])
def api_download_thread(thread_id):
    row = get_thread_row(thread_id)
    if not row:
        return jsonify({"success": False, "error": "Chat not found."}), 404
    msgs = get_message_rows(thread_id)
    lines = [
        "S.N.E.T.C.H — Document Chatbot Transcript",
        f"Chat: {row['title']}",
        f"Document: {row['doc_name'] or 'N/A'}",
        "=" * 60, "",
    ]
    for m in msgs:
        who = "You" if m["role"] == "user" else "Assistant"
        ts = datetime.fromtimestamp(m["created_at"]).strftime("%Y-%m-%d %H:%M")
        lines.append(f"[{ts}] {who}:\n{m['content']}\n")
    text = "\n".join(lines)
    buf = io.BytesIO(text.encode("utf-8"))
    buf.seek(0)
    safe_name = secure_filename(row["title"] or "chat") or "chat"
    return send_file(buf, mimetype="text/plain", as_attachment=True,
                      download_name=f"{safe_name}.txt")


@document_chatbot_bp.route("/message/<int:message_id>/feedback", methods=["POST"])
def api_message_feedback(message_id):
    data = request.get_json(force=True, silent=True) or {}
    feedback = data.get("type", "none")
    if feedback not in ("like", "dislike", "none"):
        return jsonify({"success": False, "error": "Invalid feedback type."}), 400
    set_message_feedback_row(message_id, feedback)
    return jsonify({"success": True})


def _stream_answer_response(thread_id: str, question: str, save_user_message: bool):
    """Shared SSE generator for /ask and /regenerate."""
    if thread_id not in _web_sessions:
        def err_gen():
            yield f"data: {json.dumps({'error': 'No document is loaded for this chat yet.'})}\n\n"
        return err_gen()

    if save_user_message:
        add_message_row(thread_id, "user", question)
        touch_thread(thread_id)

    def generate():
        q: queue.Queue = queue.Queue()
        state = {"text": ""}

        def on_token(t):
            state["text"] += t
            q.put(t)

        def worker():
            try:
                final = process_web_query(thread_id, question, stream=True, on_token=on_token)
                if not state["text"] and final:
                    state["text"] = final
            except Exception as e:
                q.put(("__ERROR__", str(e)))
            finally:
                q.put(None)

        threading.Thread(target=worker, daemon=True).start()

        while True:
            item = q.get()
            if item is None:
                break
            if isinstance(item, tuple) and item[0] == "__ERROR__":
                yield f"data: {json.dumps({'error': item[1]})}\n\n"
                return
            yield f"data: {json.dumps({'token': item})}\n\n"

        answer = state["text"].strip() or "This information is not available in the document."
        msg_id = add_message_row(thread_id, "assistant", answer)
        touch_thread(thread_id)
        _sync_langgraph_memory(thread_id, question, answer)
        yield f"data: {json.dumps({'done': True, 'message_id': msg_id, 'answer': answer})}\n\n"

    return generate()


@document_chatbot_bp.route("/ask", methods=["POST"])
def api_ask():
    data = request.get_json(force=True, silent=True) or {}
    thread_id = (data.get("thread_id") or "").strip()
    question = (data.get("question") or "").strip()
    if not thread_id or not get_thread_row(thread_id):
        return jsonify({"success": False, "error": "Invalid chat thread."}), 400
    if not question:
        return jsonify({"success": False, "error": "Question cannot be empty."}), 400

    return Response(
        stream_with_context(_stream_answer_response(thread_id, question, save_user_message=True)),
        mimetype="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@document_chatbot_bp.route("/regenerate", methods=["POST"])
def api_regenerate():
    data = request.get_json(force=True, silent=True) or {}
    thread_id = (data.get("thread_id") or "").strip()
    row = get_thread_row(thread_id)
    if not row:
        return jsonify({"success": False, "error": "Invalid chat thread."}), 400

    msgs = get_message_rows(thread_id)
    if not msgs or msgs[-1]["role"] != "assistant":
        return jsonify({"success": False, "error": "Nothing to regenerate yet."}), 400

    question = None
    for m in reversed(msgs):
        if m["role"] == "user":
            question = m["content"]
            break
    if not question:
        return jsonify({"success": False, "error": "No previous question found."}), 400

    delete_message_row(msgs[-1]["id"])

    return Response(
        stream_with_context(_stream_answer_response(thread_id, question, save_user_message=False)),
        mimetype="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


def register_document_chatbot(app) -> None:
    """Wire this feature's API into the main Flask app.

    document_chatbot.py deliberately does not import app.py (to avoid
    touching any other feature/file), so add these two lines once in
    app.py, near the other feature imports:

        import document_chatbot
        document_chatbot.register_document_chatbot(app)
    """
    init_document_chatbot_db()
    app.register_blueprint(document_chatbot_bp)


# Always make sure the database/upload folders exist, even if the Flask
# app never calls register_document_chatbot() (e.g. CLI-only usage).
init_document_chatbot_db()


def main():
    print("=== S.N.E.T.C.H Document AI Chatbot ===")
    print("Capabilities: summarize, key points, quiz, pages, facts, blog, flashcards,")
    print("compare documents, translate, fact-check, and 30+ more. Type 'help' for examples.\n")

    while True:
        path = input("Document path (PDF/DOCX/TXT, or Enter to skip): ").strip()
        if not path:
            print("Please provide a document path.")
            continue
        if path.lower() in ("exit", "quit"):
            return

        result = load_document(f"load {path}")
        print(f"\n{result}\n")

        doc_id = _doc_id_from_path(os.path.abspath(os.path.expanduser(path)))
        if doc_id not in _sessions:
            continue

        print("Ask anything (streaming). Commands: 'new' | 'exit' | 'help'\n")

        while True:
            question = input("You: ").strip()
            if not question:
                continue
            if question.lower() in ("exit", "quit"):
                print("Bye!")
                return
            if question.lower() == "new":
                break
            if question.lower() == "help":
                print(
                    "Examples:\n"
                    "  summarize the document\n"
                    "  main topic / key points / generate quiz\n"
                    "  what is on page 5\n"
                    "  content from page 5 to 10\n"
                    "  where is machine learning mentioned\n"
                    "  list tools and resources mentioned\n"
                    "  translate to Hindi\n"
                    "  compare C:\\docs\\a.pdf vs C:\\docs\\b.pdf\n"
                    "  fact check the main claims\n"
                )
                continue

            print("Bot: ", end="", flush=True)
            tokens = []

            def on_token(t):
                print(t, end="", flush=True)
                tokens.append(t)

            try:
                ask_document_langgraph(
                    doc_id, question, stream=True, on_token=on_token,
                )
            except Exception as e:
                print(str(e), end="")
            print("\n")


if __name__ == "__main__":
    main()
